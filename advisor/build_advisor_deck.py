#!/usr/bin/env python3
"""Build advisor/mPower_Advisor.pptx.

Source-of-truth generator for the mPower Advisor deck. Content is derived
from advisor/mpower_agentic_advisor_execution_plan.md; the evidence excerpts on
the "What It Does" slide are lifted from advisor/mpower_result_summary.html.

    uv run --with python-pptx python advisor/build_advisor_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "mPower_Advisor.pptx"

# Slide canvas (matches the original deck exactly).
SLIDE_W = 12191695
SLIDE_H = 6858000

# ---------------------------------------------------------------- design tokens
BG = RGBColor(0x0F, 0x14, 0x19)  # page background
CARD = RGBColor(0x1A, 0x20, 0x29)  # default card fill
CARD_ALT = RGBColor(0x22, 0x2A, 0x35)  # raised / emphasised card fill
BORDER = RGBColor(0x2E, 0x3A, 0x48)  # hairline border, arrow glyphs
TEXT = RGBColor(0xD7, 0xDE, 0xE7)  # primary text
MUTED = RGBColor(0x8B, 0x98, 0xA8)  # secondary text
BLUE = RGBColor(0x4F, 0xA3, 0xFF)  # accent / in-scope
GREEN = RGBColor(0x35, 0xC4, 0x7C)  # shipped value / healthy
AMBER = RGBColor(0xF2, 0xB5, 0x44)  # later / needs attention
RED = RGBColor(0xFF, 0x5F, 0x56)  # failing / over budget
INK = RGBColor(0x06, 0x10, 0x18)  # text on filled ovals
WELL = RGBColor(0x0B, 0x0F, 0x14)  # inset evidence panel (report <pre> bg)
EXCERPT = RGBColor(0xC8, 0xD4, 0xE0)  # evidence panel text (report <pre> fg)

FONT = "Segoe UI"
MONO = "Consolas"

MARGIN = 0.55  # standard content left margin (inches)
CONTENT_W = 12.22  # standard content width (inches)


# --------------------------------------------------------------------- helpers
def shape(
    slide,
    kind,
    left,
    top,
    width,
    height,
    *,
    fill=None,
    line=None,
    line_w=1.0,
    adj=None,
):
    """Add an autoshape with no text, styled from the token palette."""
    sh = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    if adj is not None:
        for i, val in enumerate(adj):
            sh.adjustments[i] = val
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.word_wrap = True
    return sh


def text(
    slide,
    left,
    top,
    width,
    height,
    paras,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=None,
):
    """Add a text box.

    ``paras`` is a list of paragraphs; each paragraph is either a run-spec dict
    or a list of run-spec dicts (for mixed styling on one line). A run spec is
    ``{"t": str, "sz": pt, "b": bool, "c": RGBColor, "f": font_name}``.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, spec in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if line_spacing:
            para.line_spacing = line_spacing
        runs = spec if isinstance(spec, list) else [spec]
        for rspec in runs:
            run = para.add_run()
            run.text = rspec["t"]
            run.font.size = Pt(rspec.get("sz", 11))
            run.font.bold = rspec.get("b", False)
            run.font.name = rspec.get("f", FONT)
            run.font.color.rgb = rspec.get("c", TEXT)
        if isinstance(spec, dict) and "space_after" in spec:
            para.space_after = Pt(spec["space_after"])
    return box


def page(prs):
    """New blank slide with the dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=BG)
    return slide


def header(slide, title, subtitle):
    """Standard slide header: accent tick, title, subtitle."""
    shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, 0.42, 0.07, 0.52, fill=BLUE)
    text(slide, 0.78, 0.36, 11.90, 0.55, [{"t": title, "sz": 26, "b": True, "c": TEXT}])
    text(slide, 0.78, 0.94, 11.90, 0.32, [{"t": subtitle, "sz": 13, "c": MUTED}])


def footer(slide, note, top=6.92):
    text(slide, MARGIN, top, CONTENT_W, 0.32, [{"t": note, "sz": 10.5, "c": MUTED}])


def card(slide, left, top, width, height, *, fill=CARD, line=BORDER, adj=0.06):
    return shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line=line,
        adj=[adj],
    )


# ------------------------------------------------------------------- slide one
def slide_title(prs):
    """Title + product usage flow illustration."""
    s = page(prs)

    text(s, 0.90, 1.00, 11.60, 0.90, [{"t": "mPower Advisor", "sz": 46, "b": True, "c": TEXT}])
    shape(s, MSO_SHAPE.RECTANGLE, 0, 2.00, 13.333, 0.03, fill=BLUE)
    text(
        s,
        0.90,
        2.24,
        11.60,
        0.42,
        [
            {
                "t": "An agentic assistant for mPower users — from “the run finished” "
                "to “I know what to do next”",
                "sz": 17.5,
                "c": BLUE,
            }
        ],
    )
    text(
        s,
        0.90,
        2.80,
        11.60,
        0.36,
        [
            {
                "t": "Automated triage of mpower.error.log, missing inputs, true warning "
                "roll-up counts and stage timing.",
                "sz": 13,
                "c": MUTED,
            }
        ],
    )

    # ---- flow illustration -------------------------------------------------
    steps = [
        {
            "kicker": "YOU RUN IT",
            "title": "Point it at a run",
            "mono": "mpower-advisor <run_dir>",
            "body": "Reads logs_<timestamp>/ in place — any finished or historical run, "
            "no tool rebuild.",
            "accent": MUTED,
        },
        {
            "kicker": "IT COLLECTS",
            "title": "Run Manifest",
            "mono": None,
            "body": "mpower.log · mpower.error.log · mpower.warning.log · mpower.lib.log · "
            "mpower.cmd · run.tcl · *.rpt",
            "accent": BLUE,
        },
        {
            "kicker": "IT ANALYZES",
            "title": "Findings",
            "mono": None,
            "body": "Rules emit severity · code · count · evidence. The model correlates and "
            "ranks — never computes.",
            "accent": BLUE,
        },
        {
            "kicker": "YOU GET",
            "title": "Run report",
            "mono": "mpower_result_summary",
            "body": ".html dashboard, or .md to paste into a ticket.",
            "accent": GREEN,
        },
    ]

    box_w, gap, box_h, box_t = 2.55, 0.44, 1.62, 3.42
    for i, st in enumerate(steps):
        x = 0.90 + i * (box_w + gap)
        card(s, x, box_t, box_w, box_h, adj=0.09)
        shape(s, MSO_SHAPE.RECTANGLE, x, box_t, 0.05, box_h, fill=st["accent"])
        text(
            s,
            x + 0.24,
            box_t + 0.17,
            box_w - 0.44,
            0.20,
            [{"t": st["kicker"], "sz": 9.5, "b": True, "c": st["accent"]}],
        )
        text(s, x + 0.24, box_t + 0.42, box_w - 0.44, 0.26, [{"t": st["title"], "sz": 13.5, "b": True, "c": TEXT}])
        body_t = box_t + 0.72
        if st["mono"]:
            text(s, x + 0.24, body_t, box_w - 0.40, 0.22, [{"t": st["mono"], "sz": 10, "c": GREEN, "f": MONO}])
            body_t += 0.28
        text(
            s,
            x + 0.24,
            body_t,
            box_w - 0.44,
            0.62,
            [{"t": st["body"], "sz": 9.5, "c": MUTED}],
            line_spacing=1.0,
        )

        if i < len(steps) - 1:
            shape(
                s,
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + 0.10,
                box_t + box_h / 2 - 0.10,
                0.24,
                0.20,
                fill=BORDER,
                adj=[0.42, 0.50],
            )

    # ---- outcome strip -----------------------------------------------------
    strip_t = box_t + box_h + 0.30
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.90,
        strip_t,
        11.53,
        0.60,
        fill=CARD_ALT,
        line=GREEN,
        line_w=0.75,
        adj=[0.30],
    )
    text(
        s,
        1.18,
        strip_t + 0.19,
        11.00,
        0.26,
        [
            [
                {"t": "OUTCOME   ", "sz": 10, "b": True, "c": GREEN},
                {"t": "A ranked ", "sz": 11.5, "c": TEXT},
                {"t": "“Recommended Next Actions”", "sz": 11.5, "b": True, "c": TEXT},
                {"t": " list — every number cited back to a file and line.", "sz": 11.5, "c": TEXT},
            ]
        ],
    )

    footer(s, "Draft for review  ·  July 2026  ·  Siemens EDA — mPower", top=6.60)


# ------------------------------------------------------------------- slide two
def slide_what_it_does(prs):
    """Four in-scope capabilities, each shown with real reference-run output."""
    s = page(prs)
    header(
        s,
        "What It Does",
        "Four capabilities in the first release — each shown with real output from the "
        "BRCM rts_top reference run",
    )

    # Excerpts are lifted from advisor/mpower_result_summary.html. Mono runs are
    # space-padded so columns line up within a line.
    caps = [
        {
            "n": "1",
            "colour": GREEN,
            "title": "Run digest / dashboard",
            "q": "“Did my run do what I asked, and can I trust the result?”",
            "evidence": [
                [{"t": "Run outcome       ", "c": MUTED}, {"t": "Completed · final.db saved", "c": GREEN}],
                [{"t": "Wall / peak mem   ", "c": MUTED}, {"t": "11h 18m  ·  98.7 GB", "c": EXCERPT}],
                [
                    {"t": "Worst dynamic IR  ", "c": MUTED},
                    {"t": "14.7%", "c": RED},
                    {"t": "  = 111.6 mV of 0.76 V", "c": EXCERPT},
                ],
                [
                    {"t": "EM violations     ", "c": MUTED},
                    {"t": "0", "c": GREEN},
                    {"t": "      Hard errors  ", "c": MUTED},
                    {"t": "4", "c": AMBER},
                ],
            ],
        },
        {
            "n": "2",
            "colour": GREEN,
            "title": "Input & setup validation",
            "q": "“What did mPower silently work around?”",
            "evidence": [
                [
                    {"t": "PWR-327  ", "c": MUTED},
                    {"t": "523,976", "c": AMBER},
                    {"t": " of 4,045,248 nets: no arrival time", "c": EXCERPT},
                ],
                [{"t": "PWR-297  ", "c": MUTED}, {"t": "Timescale is not defined in the TWF", "c": EXCERPT}],
                [{"t": "warning  ", "c": MUTED}, {"t": "No instances found in twf", "c": EXCERPT}],
                [
                    {"t": "→ ", "c": AMBER},
                    {"t": "4,888", "c": AMBER},
                    {"t": " nets fell back to toggle_rate 0.07", "c": EXCERPT},
                ],
            ],
        },
        {
            "n": "3",
            "colour": BLUE,
            "title": "Error triage",
            "q": "“It failed or looks wrong — why, and what do I change?”",
            "evidence": [
                [{"t": "Error: Clock ...u_rts_clk_rst_ctrl/dnt_ckbuf_rts_clk/o", "c": EXCERPT}],
                [
                    {"t": "       already exist, ", "c": EXCERPT},
                    {"t": "command ignored", "c": RED},
                    {"t": "     × 4", "c": MUTED},
                ],
                [{"t": "→ mPower kept the first definition — on the same", "c": BLUE}],
                [{"t": "  block that has the worst dynamic IR (14.7%).", "c": BLUE}],
            ],
        },
        {
            "n": "4",
            "colour": BLUE,
            "title": "Performance advisor",
            "q": "“It took 11 hours. Where did the time go?”",
            "evidence": [
                [
                    {"t": "Dynamic IR 50 ns @ 5 ps  ", "c": MUTED},
                    {"t": "7h 45m", "c": EXCERPT},
                    {"t": "   69%", "c": RED},
                ],
                [
                    {"t": "SPEF + TWF + Power       ", "c": MUTED},
                    {"t": "2h 08m", "c": EXCERPT},
                    {"t": "   19%", "c": AMBER},
                ],
                [{"t": "→ 6.3h CPU vs 7.8h wall: not parallelizing well", "c": BLUE}],
                [{"t": "→ distributed power off · 16 of 128 slots used", "c": BLUE}],
            ],
        },
    ]

    col_w, col_gap = 6.00, 0.23
    row_h, row_gap = 2.10, 0.12
    for i, cap in enumerate(caps):
        x = MARGIN + (i % 2) * (col_w + col_gap)
        y = 1.42 + (i // 2) * (row_h + row_gap)
        card(s, x, y, col_w, row_h)
        ov = shape(s, MSO_SHAPE.OVAL, x + 0.24, y + 0.20, 0.34, 0.34, fill=cap["colour"])
        ov.text_frame.word_wrap = False
        text(
            s,
            x + 0.24,
            y + 0.25,
            0.34,
            0.24,
            [{"t": cap["n"], "sz": 12, "b": True, "c": INK}],
            align=PP_ALIGN.CENTER,
        )
        text(s, x + 0.72, y + 0.22, col_w - 0.96, 0.30, [{"t": cap["title"], "sz": 15, "b": True, "c": TEXT}])
        text(s, x + 0.26, y + 0.68, col_w - 0.52, 0.28, [{"t": cap["q"], "sz": 12, "c": cap["colour"]}])

        # evidence well — mirrors the report's <pre> styling
        shape(
            s,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + 0.26,
            y + 1.02,
            col_w - 0.52,
            0.88,
            fill=WELL,
            line=BORDER,
            line_w=0.75,
            adj=[0.06],
        )
        lines = [[dict(r, sz=8.5, f=MONO) for r in line] for line in cap["evidence"]]
        text(s, x + 0.42, y + 1.13, col_w - 0.82, 0.68, lines, line_spacing=1.18)

    # ---- coming next -------------------------------------------------------
    strip_t = 1.42 + 2 * row_h + row_gap + 0.14
    shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN,
        strip_t,
        CONTENT_W,
        0.88,
        fill=None,
        line=AMBER,
        line_w=1.0,
        adj=[0.18],
    )
    text(s, 0.84, strip_t + 0.18, 1.45, 0.24, [{"t": "COMING NEXT", "sz": 10.5, "b": True, "c": AMBER}])
    text(s, 0.84, strip_t + 0.44, 1.45, 0.24, [{"t": "Phases 4–5", "sz": 9.5, "c": MUTED}])

    nxt = [
        (
            "EM / IR-drop debug",
            "“This instance violates — what causes it and what’s the fix?”",
            "Needs hotspot ↔ power-density ↔ grid-R ↔ via-density correlation.",
        ),
        (
            "Setup assistant",
            "“How do I set this flow up correctly in the first place?”",
            "Needs a canonical flow model mined from the 826 regression runScripts.",
        ),
    ]
    for i, (title, question, needs) in enumerate(nxt):
        x = 2.42 + i * 5.25
        text(s, x, strip_t + 0.13, 5.00, 0.24, [{"t": title, "sz": 12.5, "b": True, "c": TEXT}])
        text(s, x, strip_t + 0.38, 5.00, 0.22, [{"t": question, "sz": 9.5, "c": AMBER}])
        text(s, x, strip_t + 0.59, 5.00, 0.22, [{"t": needs, "sz": 9.5, "c": MUTED}])

    footer(
        s,
        "Delivered as mpower_result_summary.html (dashboard) and .md (paste into a ticket).",
        top=strip_t + 1.02,
    )


# ----------------------------------------------------------------- slide three
def slide_architecture(prs):
    s = page(prs)
    header(s, "Architecture", "Five layers — each independently useful and independently testable")

    layers = [
        (
            "INPUT",
            MUTED,
            "Run directory",
            "mpower.log · mpower.error.log · mpower.warning.log · mpower.lib.log · "
            "mpower.cmd · run.tcl · *.rpt",
            CARD_ALT,
        ),
        (
            "L1",
            BLUE,
            "Collectors — discovery & normalization",
            "Emits the Run Manifest: one versioned JSON that everything downstream consumes",
            CARD,
        ),
        (
            "L2",
            BLUE,
            "Knowledge base",
            "1,767 .mhelp codes · 390 ghelp/*.help commands · 826 regression runScripts · playbooks",
            CARD,
        ),
        (
            "L3",
            BLUE,
            "Analyzers — deterministic rules → Findings",
            "completeness · error triage · warning roll-up · performance · margin · config sanity",
            CARD,
        ),
        (
            "L4",
            BLUE,
            "Reasoning — LLM correlation & narrative",
            "Cross-domain correlation, impact ranking, prose · + agentic tool-use (Phase 3+)",
            CARD,
        ),
        (
            "L5",
            BLUE,
            "Presentation",
            "mpower_result_summary.html · mpower_result_summary.md · interactive chat",
            CARD,
        ),
    ]

    y = 1.46
    for i, (tag, colour, title, body, fill) in enumerate(layers):
        card(s, MARGIN, y, 9.15, 0.73, fill=fill, adj=0.12)
        text(s, 0.77, y + 0.06, 0.62, 0.61, [{"t": tag, "sz": 13, "b": True, "c": colour}], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.47, y + 0.11, 7.95, 0.28, [{"t": title, "sz": 13, "b": True, "c": TEXT}])
        text(s, 1.47, y + 0.42, 7.95, 0.26, [{"t": body, "sz": 10.2, "c": MUTED}])
        if i < len(layers) - 1:
            shape(s, MSO_SHAPE.DOWN_ARROW, 0.97, y + 0.74, 0.16, 0.11, fill=BORDER, adj=[0.5, 0.5])
        y += 0.86

    # ---- design principles -------------------------------------------------
    card(s, 9.92, 1.46, 2.85, 5.19, adj=0.05)
    text(s, 10.16, 1.66, 2.37, 0.30, [{"t": "DESIGN PRINCIPLES", "sz": 10.5, "b": True, "c": BLUE}])

    # Explicit tops: block heights differ (title and body wrap by different
    # amounts), so a uniform pitch would leave visually uneven gaps.
    principles = [
        (
            2.06,
            "Deterministic first, LLM second",
            "Every number comes from a parser. The model explains and ranks — never computes.",
        ),
        (3.36, "Cite everything", "Every finding links to source file and line — auditable in one click."),
        (
            4.47,
            "Degrade gracefully",
            "Crashed runs, partial runs, and runs with no report_time calls still report.",
        ),
        (5.78, "Offline-capable", "Usable with no external network — and with no LLM at all."),
    ]
    for py, title, body in principles:
        text(
            s,
            10.16,
            py,
            2.37,
            1.10,
            [
                {"t": title, "sz": 11.5, "b": True, "c": TEXT, "space_after": 3},
                {"t": body, "sz": 9.8, "c": MUTED},
            ],
            line_spacing=1.0,
        )

    footer(
        s,
        "Key decision: the Run Manifest decouples parsing from analysis — the whole system is "
        "testable without a single LLM call.",
    )


# ------------------------------------------------------------------ slide four
def slide_roadmap(prs):
    s = page(prs)
    header(s, "Roadmap", "What the user gets, in the order it ships")

    milestones = [
        (
            "0",
            GREEN,
            "Groundwork",
            "Phase 0",
            "Labeled corpus of 20–30 historical runs, the Run Manifest schema, a normalized "
            ".mhelp <type> field, and the eval harness.",
            "2 EW",
        ),
        (
            "1",
            GREEN,
            "A run report you can trust",
            "Phase 1",
            "Deterministic digest — KPIs, stage timing, max-of-roll-up warning counts, missing "
            "inputs. No LLM. Shippable to AEs on its own.",
            "2–3 EW",
        ),
        (
            "2",
            BLUE,
            "It explains itself",
            "Phase 2",
            "Cross-domain correlation, severity ranking and a prioritized next-actions list — "
            "strictly grounded, every claim cited.",
            "2–3 EW",
        ),
        (
            "3",
            BLUE,
            "You can ask it questions",
            "Phase 3",
            "Conversational follow-up (“show me the nets without_delay”), tool-use over the "
            "manifest, multi-run diff vs a prior run.",
            "2–3 EW",
        ),
        (
            "4",
            AMBER,
            "It finds the root cause",
            "Phase 4",
            "EM / IR-drop debug — hotspot clustering, blame attribution, decap adequacy, "
            "concrete mitigation proposals.",
            "10–14 EW",
        ),
        (
            "5",
            AMBER,
            "It catches problems before you run",
            "Phase 5",
            "Setup assistant — pre-flight flow validation, guided setup, an analyze_run Tcl "
            "command, GUI2 panel / CI gating.",
            "6–10 EW",
        ),
    ]

    row_top, row_h, row_pitch = 1.44, 0.66, 0.765
    shape(
        s,
        MSO_SHAPE.RECTANGLE,
        0.86,
        row_top + 0.10,
        0.04,
        row_pitch * (len(milestones) - 1) + 0.30,
        fill=BORDER,
    )

    for i, (num, colour, title, phase, body, ew) in enumerate(milestones):
        y = row_top + i * row_pitch
        card(s, 1.42, y, 11.35, row_h, adj=0.14)
        ov = shape(s, MSO_SHAPE.OVAL, 0.71, y + 0.17, 0.33, 0.33, fill=colour)
        ov.text_frame.word_wrap = False
        text(
            s,
            0.71,
            y + 0.22,
            0.33,
            0.24,
            [{"t": num, "sz": 11.5, "b": True, "c": INK}],
            align=PP_ALIGN.CENTER,
        )
        text(
            s,
            1.68,
            y + 0.08,
            8.90,
            0.28,
            [
                [
                    {"t": title, "sz": 12.5, "b": True, "c": TEXT},
                    {"t": "     " + phase, "sz": 10, "c": MUTED},
                ]
            ],
        )
        text(s, 1.68, y + 0.37, 9.40, 0.26, [{"t": body, "sz": 10.4, "c": MUTED}])
        shape(
            s,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            11.32,
            y + 0.16,
            1.22,
            0.34,
            fill=CARD_ALT,
            line=colour,
            line_w=0.75,
            adj=[0.4],
        )
        text(
            s,
            11.32,
            y + 0.22,
            1.22,
            0.24,
            [{"t": ew, "sz": 11, "b": True, "c": colour}],
            align=PP_ALIGN.CENTER,
        )

    # Runs alongside every phase, and is staffed differently — hence the separate band.
    band_t = row_top + len(milestones) * row_pitch + 0.02
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.42, band_t, 11.35, 0.58, fill=None, line=BLUE, line_w=1.0, adj=[0.26])
    text(
        s,
        1.68,
        band_t + 0.10,
        10.85,
        0.40,
        [
            [
                {
                    "t": "Runs alongside every phase — remediation knowledge  "
                    "(4–6 expert-weeks, not dev time):  ",
                    "sz": 11,
                    "b": True,
                    "c": BLUE,
                },
                {
                    "t": "976 of the 1,767 .mhelp codes carry <action> No_Information. "
                    "Filling <action> for the top codes by frequency turns a finding into a fix.",
                    "sz": 11,
                    "c": MUTED,
                },
            ]
        ],
        line_spacing=1.0,
    )

    footer(
        s,
        "First shippable value: ~4–5 EW (Phases 0+1)   ·   Full vision: ~24–35 dev-EW "
        "≈ 2 engineers × 3–4 months, plus 4–6 domain-expert weeks",
        top=band_t + 0.66,
    )


# --------------------------------------------------------------------- driver
def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    slide_title(prs)
    slide_what_it_does(prs)
    slide_architecture(prs)
    slide_roadmap(prs)

    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
